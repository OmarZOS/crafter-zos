
import collections 
import collections.abc
from pptx import Presentation



# data = {
#     "2221":"Nobody",
#     "2222":"keiner",
#     "2223":"Niemand",
#     "2224":"Nemo",
#     "2225":"Outis",
#     "2226":"Personne"
# }



def search_and_replace(data, input, output):
    """"search and replace text in PowerPoint while preserving formatting"""
    #Useful Links ;)
    #https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
    #https://stackoverflow.com/questions/45247042/how-to-keep-original-text-formatting-of-text-with-python-powerpoint
    from pptx import Presentation
    prs = Presentation(input)
    for search_str,repl_str in data.items():
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if(shape.text.find(search_str))!=-1:
                        text_frame = shape.text_frame
                        cur_text = text_frame.paragraphs[0].runs[0].text
                        new_text = cur_text.replace(str(search_str), str(repl_str))
                        text_frame.paragraphs[0].runs[0].text = new_text
    prs.save(output)


# search_and_replace(data,"diplome.pptx","dip2.pptx")